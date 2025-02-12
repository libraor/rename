import os
import tkinter as tk
from tkinter import filedialog, messagebox
import traceback
from tkinter.ttk import Progressbar
import logging
from openai import OpenAI
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
import fitz  # PyMuPDF
from PIL import Image
from concurrent.futures import ThreadPoolExecutor, as_completed
from paddleocr import PaddleOCR
import numpy as np
import time  # 导入 time 模块用于计时
from skimage.metrics import structural_similarity as ssim
import cv2

# 从 pw.py 文件中导入配置信息
from pw import API_KEY, BASE_URL, MODEL_NAME

os.environ['LIBPNG_WARNING_LEVEL'] = '2'  # 设置为2可以抑制警告
logging.basicConfig(filename='app.log', level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')

# 更改 API 密钥和基础 URL 为火山接口的信息
client = OpenAI(
    api_key=API_KEY,  # 火山接口 API 密钥
    base_url=BASE_URL,  # 火山接口域名
)

def extract_time_openai(text):
    """
    使用火山接口模型从文本中提取时间信息
    """
    try:
        response = client.chat.completions.create(
            model=MODEL_NAME,  # 火山接口模型名称
            messages=[{
                "role": "user",
                "content": f"假设你是文件重命名助手，分析文件生成时间与主要内容，以 “yyyymmdd_标题” 格式返回。若无法识别时间，以 “00000000_标题” 格式输出，标题简洁，不超 20 字。不需要任何解释。不需要解析过程。{text}"
            }]
        )
        logging.info("成功调用火山接口 API")
        result = response.choices[0].message.content.strip()
        return result
    except Exception as e:
        error_message = f"调用火山接口 API 时出错：{e}"
        logging.error(error_message, exc_info=True)
        messagebox.showerror("错误", error_message)
        return None

def read_docx(file_path):
    """
    读取docx文件内容，包括段落和表格
    """
    try:
        with open(file_path, 'rb') as f:
            doc = Document(f)
        
        full_text = []

        # 读取段落内容
        for para in doc.paragraphs:
            full_text.append(para.text)

        # 读取表格内容
        for table in doc.tables:
            table_text = []
            for row in table.rows:
                row_text = []
                for cell in row.cells:
                    row_text.append(cell.text)
                table_text.append('\t'.join(row_text))
            full_text.append('\n'.join(table_text))

        content = '\n'.join(full_text)
        print(f"成功读取 {file_path} 内容: {content[:100]}...")  # 直接在终端中输出前100个字符以避免输出过长
        return content
    except Exception as e:
        error_message = f"读取 {file_path} 时出错：{e}"
        logging.error(error_message, exc_info=True)
        messagebox.showerror("错误", error_message)
        return None

def read_xlsx(file_path):
    """
    读取xlsx文件内容
    """
    try:
        wb = load_workbook(file_path, read_only=True)
        full_text = []
        for sheet in wb.sheetnames:
            ws = wb[sheet]
            for row in ws.iter_rows(values_only=True):
                valid_cells = [str(cell) for cell in row if cell is not None]
                full_text.append('\t'.join(valid_cells))
        return '\n'.join(full_text)
    except Exception as e:
        logging.error(f"读取 {file_path} 时出错：{e}", exc_info=True)
        return None

def read_pptx(file_path):
    """
    读取pptx文件内容
    """
    try:
        with open(file_path, 'rb') as f:
            prs = Presentation(f)
            full_text = [shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")]
        return '\n'.join(full_text)
    except Exception as e:
        logging.error(f"读取 {file_path} 时出错：{e}", exc_info=True)
        return None

# 初始化 PaddleOCR
ocr = PaddleOCR(use_angle_cls=True, lang='ch')  # 初始化 PaddleOCR，使用中文

def read_pdf(file_path):
    try:
        doc = fitz.open(file_path)
        full_text = []
        for page_num in range(len(doc)):
            page = doc.load_page(page_num)
            text = page.get_text()
            if text.strip():  # 如果页面有可提取的文本
                full_text.append(text)
            else:  # 如果页面没有可提取的文本，尝试使用 OCR
                logging.info(f"第 {page_num + 1} 页未提取到文本，尝试使用 PaddleOCR")
                pix = page.get_pixmap()
                img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                # 将 PIL 图像转换为 numpy 数组
                img_array = np.array(img)
                try:
                    ocr_result = ocr.ocr(img_array, cls=True)  # 使用 PaddleOCR 进行 OCR
                except Exception as e:
                    logging.error(f"PaddleOCR 处理第 {page_num + 1} 页时出错：{e}", exc_info=True)
                    continue
                if not ocr_result:  # 检查 OCR 结果是否为空
                    logging.error("OCR 结果为空")
                    continue
                ocr_text = []
                for line in ocr_result:
                    for word in line:
                        ocr_text.append(word[1][0])  # 提取识别的文本
                full_text.append(' '.join(ocr_text))
                logging.info(f"第 {page_num + 1} 页 OCR 提取完成")
        return '\n'.join(full_text)
    except Exception as e:
        logging.error(f"读取 {file_path} 时出错：{e}", exc_info=True)
        return None

def read_image(file_path):
    """
    使用 PaddleOCR 读取图片文件内容
    """
    try:
        img = Image.open(file_path)
        # 将 PIL 图像转换为 numpy 数组
        img_array = np.array(img)
        try:
            ocr_result = ocr.ocr(img_array, cls=True)  # 使用 PaddleOCR 进行 OCR
        except Exception as e:
            logging.error(f"PaddleOCR 处理图片 {file_path} 时出错：{e}", exc_info=True)
            return None
        if not ocr_result:  # 检查 OCR 结果是否为空
            logging.error("OCR 结果为空")
            return None
        ocr_text = []
        for line in ocr_result:
            for word in line:
                ocr_text.append(word[1][0])  # 提取识别的文本
        return ' '.join(ocr_text)
    except Exception as e:
        logging.error(f"读取图片 {file_path} 时出错：{e}", exc_info=True)
        return None

def read_text_file(file_path):
    """
    尝试使用不同编码读取文本文件
    """
    encodings = ['utf-8', 'gbk', 'latin1']
    for encoding in encodings:
        try:
            with open(file_path, 'r', encoding=encoding) as f:
                return f.read()
        except UnicodeDecodeError:
            continue
    logging.error(f"无法解码文件 {file_path}，请检查文件编码。")
    messagebox.showerror("错误", f"无法解码文件 {file_path}，请检查文件编码。")
    return None

def get_files(directory):
    """
    让用户选择一个文件夹，并返回该文件夹下的所有文件路径列表
    """
    if directory:
        file_list = []
        for root, dirs, files in os.walk(directory):
            for file in files:
                file_path = os.path.join(root, file)
                file_list.append(file_path)
        return file_list
    return []

def get_file_content(file_path):
    """
    根据文件扩展名选择读取函数
    """
    file_ext = os.path.splitext(file_path)[1].lower()
    if file_ext == '.docx':
        return read_docx(file_path)
    elif file_ext == '.xlsx':
        return read_xlsx(file_path)
    elif file_ext == '.pptx':
        return read_pptx(file_path)
    elif file_ext == '.pdf':
        return read_pdf(file_path)
    elif file_ext in ['.jpg', '.jpeg', '.png', '.bmp', '.gif']:
        return read_image(file_path)
    else:
        return read_text_file(file_path)

def sanitize_filename(filename):
    """
    过滤掉文件名中的非法字符
    """
    return "".join(x for x in filename if x.isalnum() or x in "._- ")

def process_file(file, progress_bar, total_files):
    """
    处理单个文件
    """
    try:
        start_time = time.time()  # 记录开始时间
        content = get_file_content(file)
        if content is None:
            print(f"读取文件 {file} 失败")  # 输出读取文件失败信息
            return (file, None, None)

        content_length = len(content)  # 记录文件内容长度

        time_info = extract_time_openai(content)
        if time_info:
            file_dir = os.path.dirname(file)
            file_ext = os.path.splitext(file)[1]
            new_file_name = sanitize_filename(f"{time_info}{file_ext}")
            new_file_path = os.path.join(file_dir, new_file_name)
            os.rename(file, new_file_path)
            end_time = time.time()  # 记录结束时间
            elapsed_time = end_time - start_time  # 计算用时
            print(f"文件 {file} 已重命名为 {new_file_path}，用时: {elapsed_time:.2f} 秒")  # 输出成功信息和用时
            logging.info(f"文件 {file} 已重命名为 {new_file_path}")
            return (file, elapsed_time, content_length)
        else:
            print(f"文件 {file} 处理失败，未获取到时间信息")  # 输出处理失败信息
            return (file, None, content_length)
    except Exception as e:
        logging.error(f"处理文件 {file} 时出错：{e}", exc_info=True)
        error_message = f"处理文件 {file} 时出错：{e}\n详细信息：\n{traceback.format_exc()}"
        messagebox.showerror("错误", error_message)
        print(f"处理文件 {file} 失败")  # 输出处理文件失败信息
        return (file, None, None)

def split_pdf_by_layout(pdf_path, output_dir):
    """
    根据布局分割PDF文件
    """
    import fitz

    def analyze_layout(page):
        layout_features = []
        text_dict = page.get_text("dict")
        if "blocks" in text_dict:
            for block in text_dict["blocks"]:
                if "lines" in block:
                    for line in block["lines"]:
                        for span in line["spans"]:
                            layout_features.append((span["size"], span["font"]))
        return layout_features

    def calculate_image_similarity(img1, img2):
        # 将图像转换为灰度图像
        gray1 = cv2.cvtColor(img1, cv2.COLOR_BGR2GRAY)
        gray2 = cv2.cvtColor(img2, cv2.COLOR_BGR2GRAY)
        # 计算结构相似性
        similarity, _ = ssim(gray1, gray2, full=True)
        return similarity

    try:
        doc = fitz.open(pdf_path)
        layout_changes = []
        prev_layout = None
        prev_image = None

        for page_num in range(doc.page_count):
            page = doc[page_num]
            current_layout = analyze_layout(page)
            pix = page.get_pixmap()
            img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
            img_array = np.array(img)

            if prev_layout is None:
                prev_layout = current_layout
                prev_image = img_array
                continue

            # 检查文本布局变化
            if abs(len(current_layout) - len(prev_layout)) > 10:
                layout_changes.append(page_num)

            # 检查图像相似度
            if prev_image is not None:
                similarity = calculate_image_similarity(prev_image, img_array)
                if similarity < 0.9:  # 如果相似度低于 0.9，则认为是新的布局
                    layout_changes.append(page_num)

            prev_layout = current_layout
            prev_image = img_array

        current_file_pages = []
        file_index = 0
        for i in range(doc.page_count):
            current_file_pages.append(i)
            if i in layout_changes or i == doc.page_count - 1:
                new_doc = fitz.open()
                for num in current_file_pages:
                    new_doc.insert_pdf(doc, from_page=num, to_page=num)
                output_path = os.path.join(output_dir, f'split_{file_index}.pdf')
                new_doc.save(output_path)
                logging.info(f"分割后的 PDF 文件已保存到 {output_path}")
                file_index += 1
                current_file_pages = []

        # 关闭原始文档
        doc.close()

        # 删除原始 PDF 文件
        try:
            os.remove(pdf_path)
            logging.info(f"原始 PDF 文件 {pdf_path} 已删除")
        except Exception as e:
            error_message = f"删除原始 PDF 文件 {pdf_path} 时出错：{e}"
            logging.error(error_message, exc_info=True)
            messagebox.showerror("错误", error_message)
    except Exception as e:
        error_message = f"分割 PDF 文件 {pdf_path} 时出错：{e}"
        logging.error(error_message, exc_info=True)
        messagebox.showerror("错误", error_message)

def process_files():
    """
    处理多个文件
    """
    directory = filedialog.askdirectory()
    if not directory:
        return

    # 获取目录中的所有文件
    files = get_files(directory)

    # 处理 PDF 文件
    pdf_files = [file for file in files if file.lower().endswith('.pdf')]
    for pdf_file in pdf_files:
        logging.info(f"开始分割 PDF 文件: {pdf_file}")
        split_pdf_by_layout(pdf_file, directory)
        logging.info(f"完成分割 PDF 文件: {pdf_file}")

    # 重新获取目录中的所有文件
    files = get_files(directory)

    if not files:
        return

    total_files = len(files)
    progress_bar = Progressbar(root, orient='horizontal', length=300, mode='determinate')
    progress_bar.pack(pady=10)
    progress_bar['maximum'] = total_files
    processed_count = 0

    def update_progress(value):
        percent = int((value / total_files) * 100)
        progress_label.config(text=f"{percent}% 完成")
        progress_bar['value'] = value
        root.update_idletasks()

    file_times = {}  # 用于存储每个文件的处理时间
    file_sizes = {}  # 用于存储每个文件的内容长度

    total_elapsed_time = 0  # 总处理时间
    total_content_length = 0  # 总内容长度

    with ThreadPoolExecutor(max_workers=1) as executor:  # 创建一个线程池，最大线程数为 1
        futures = {executor.submit(process_file, file, progress_bar, total_files): file for file in files}
        for future in as_completed(futures):
            try:
                file, elapsed_time, content_length = future.result()
                if elapsed_time is not None:
                    file_times[file] = elapsed_time
                    total_elapsed_time += elapsed_time
                if content_length is not None:
                    file_sizes[file] = content_length
                    total_content_length += content_length
                processed_count += 1
                update_progress(processed_count)
            except Exception as e:
                logging.error(f"处理文件时出错：{e}")

    # 输出每个文件的处理时间
    print("\n每个文件的处理时间：")
    for file, elapsed_time in file_times.items():
        print(f"{file}: {elapsed_time:.2f} 秒")

    # 输出每个文件的内容长度
    print("\n每个文件的内容长度：")
    for file, content_length in file_sizes.items():
        print(f"{file}: {content_length} 字节")

    # 输出总的文件处理时间和总的文件内容长度
    print(f"\n总的文件处理时间: {total_elapsed_time:.2f} 秒")
    print(f"总的文件内容长度: {total_content_length} 字节")

    messagebox.showinfo("完成", "文件处理已完成。")
    progress_bar.destroy()  # 关闭进度条
    progress_label.config(text="处理已完成")

# 创建主窗口
root = tk.Tk()
root.title("文件批量处理工具")

progress_label = tk.Label(root, text="0% 完成")
progress_label.pack(pady=5)

# 创建一个按钮，点击时调用 process_files 函数
process_button = tk.Button(root, text="选择文件夹并处理文件", command=process_files)
process_button.pack(pady=20)

# 启动主事件循环，使窗口保持显示并响应用户操作
root.mainloop()