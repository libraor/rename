# file_reader.py
import os
import fitz  # PyMuPDF
from PIL import Image
import numpy as np
import cv2
from skimage.metrics import structural_similarity as ssim
from paddleocr import PaddleOCR
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
import logging

# 初始化 PaddleOCR
ocr = PaddleOCR(use_angle_cls=True, lang='ch')  # 初始化 PaddleOCR，使用中文

def read_docx(file_path):
    """
    读取docx文件内容，包括段落和表格
    """
    try:
        with open(file_path, 'rb') as f:
            doc = Document(f)
        
        full_text = []

        # 读取段落内容
        for para in doc.paragraphs:
            full_text.append(para.text)

        # 读取表格内容
        for table in doc.tables:
            table_text = []
            for row in table.rows:
                row_text = []
                for cell in row.cells:
                    row_text.append(cell.text)
                table_text.append('\t'.join(row_text))
            full_text.append('\n'.join(table_text))

        content = '\n'.join(full_text)
        print(f"成功读取 {file_path} 内容: {content[:100]}...")  # 直接在终端中输出前100个字符以避免输出过长
        return content
    except Exception as e:
        logging.error(f"读取 {file_path} 时出错：{e}", exc_info=True)
        return None

def read_xlsx(file_path):
    """
    读取xlsx文件内容
    """
    try:
        wb = load_workbook(file_path, read_only=True)
        full_text = []
        for sheet in wb.sheetnames:
            ws = wb[sheet]
            for row in ws.iter_rows(values_only=True):
                valid_cells = [str(cell) for cell in row if cell is not None]
                full_text.append('\t'.join(valid_cells))
        return '\n'.join(full_text)
    except Exception as e:
        logging.error(f"读取 {file_path} 时出错：{e}", exc_info=True)
        return None

def read_pptx(file_path):
    """
    读取pptx文件内容
    """
    try:
        with open(file_path, 'rb') as f:
            prs = Presentation(f)
            full_text = [shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")]
        return '\n'.join(full_text)
    except Exception as e:
        logging.error(f"读取 {file_path} 时出错：{e}", exc_info=True)
        return None

def read_pdf(file_path):
    try:
        doc = fitz.open(file_path)
        full_text = []
        for page_num in range(len(doc)):
            page = doc.load_page(page_num)
            text = page.get_text()
            if text.strip():  # 如果页面有可提取的文本
                full_text.append(text)
            else:  # 如果页面没有可提取的文本，尝试使用 OCR
                logging.info(f"第 {page_num + 1} 页未提取到文本，尝试使用 PaddleOCR")
                pix = page.get_pixmap()
                img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                # 将 PIL 图像转换为 numpy 数组
                img_array = np.array(img)
                try:
                    ocr_result = ocr.ocr(img_array, cls=True)  # 使用 PaddleOCR 进行 OCR
                except Exception as e:
                    logging.error(f"PaddleOCR 处理第 {page_num + 1} 页时出错：{e}", exc_info=True)
                    continue
                if not ocr_result:  # 检查 OCR 结果是否为空
                    logging.error("OCR 结果为空")
                    continue
                ocr_text = []
                for line in ocr_result:
                    for word in line:
                        ocr_text.append(word[1][0])  # 提取识别的文本
                full_text.append(' '.join(ocr_text))
                logging.info(f"第 {page_num + 1} 页 OCR 提取完成")
        return '\n'.join(full_text)
    except Exception as e:
        logging.error(f"读取 {file_path} 时出错：{e}", exc_info=True)
        return None

def read_image(file_path):
    """
    使用 PaddleOCR 读取图片文件内容
    """
    try:
        img = Image.open(file_path)
        # 将 PIL 图像转换为 numpy 数组
        img_array = np.array(img)
        try:
            ocr_result = ocr.ocr(img_array, cls=True)  # 使用 PaddleOCR 进行 OCR
        except Exception as e:
            logging.error(f"PaddleOCR 处理图片 {file_path} 时出错：{e}", exc_info=True)
            return None
        if not ocr_result:  # 检查 OCR 结果是否为空
            logging.error("OCR 结果为空")
            return None
        ocr_text = []
        for line in ocr_result:
            for word in line:
                ocr_text.append(word[1][0])  # 提取识别的文本
        return ' '.join(ocr_text)
    except Exception as e:
        logging.error(f"读取图片 {file_path} 时出错：{e}", exc_info=True)
        return None

def read_text_file(file_path):
    """
    尝试使用不同编码读取文本文件
    """
    encodings = ['utf-8', 'gbk', 'latin1']
    for encoding in encodings:
        try:
            with open(file_path, 'r', encoding=encoding) as f:
                return f.read()
        except UnicodeDecodeError:
            continue
    logging.error(f"无法解码文件 {file_path}，请检查文件编码。")
    return None

def get_file_content(file_path):
    """
    根据文件扩展名选择读取函数
    """
    file_ext = os.path.splitext(file_path)[1].lower()
    if file_ext == '.docx':
        return read_docx(file_path)
    elif file_ext == '.xlsx':
        return read_xlsx(file_path)
    elif file_ext == '.pptx':
        return read_pptx(file_path)
    elif file_ext == '.pdf':
        return read_pdf(file_path)
    elif file_ext in ['.jpg', '.jpeg', '.png', '.bmp', '.gif']:
        return read_image(file_path)
    else:
        return read_text_file(file_path)